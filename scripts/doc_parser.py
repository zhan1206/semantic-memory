#!/usr/bin/env python3
"""
from __future__ import annotations
Semantic Memory — 文档解析器 v2
支持 PDF / TXT / Markdown / DOCX / XLSX / PPTX 自动解析、分块、表格检测

增强功能:
- PDF 多层提取（文字 + 表格结构 + 图片说明）
- 扫描 PDF OCR 提示
- 表格智能识别与 Markdown 格式化
- XLSX / PPTX 支持
"""
import os
import re
import logging

logger = logging.getLogger("doc_parser")

# ─── 入口 ──────────────────────────────────────────────────

def parse_file(file_path: str) -> str:
    """
    解析文档文件，返回纯文本内容
    支持格式及优先级:
      - .txt, .md, .markdown  → 文本
      - .pdf                   → PDF（含表格检测）
      - .docx                  → Word（含表格）
      - .xlsx                  → Excel（每个 Sheet 一节）
      - .pptx                  → PPT（每张幻灯片一段）
    """
    ext = os.path.splitext(file_path)[1].lower()

    parsers = {
        ".txt": _parse_text,
        ".md": _parse_text,
        ".markdown": _parse_text,
        ".pdf": _parse_pdf,
        ".docx": _parse_docx,
        ".doc": _parse_docx,
        ".xlsx": _parse_xlsx,
        ".pptx": _parse_pptx,
    }

    parser = parsers.get(ext)
    if parser is None:
        raise ValueError(
            f"不支持的文件格式: {ext}（支持 .txt/.md/.pdf/.docx/.xlsx/.pptx）"
        )
    return parser(file_path)


# ─── 文本解析 ──────────────────────────────────────────────

def _parse_text(file_path: str) -> str:
    """解析纯文本 / Markdown 文件，自动检测编码"""
    import chardet

    with open(file_path, "rb") as f:
        raw = f.read()

    detected = chardet.detect(raw)
    encoding = detected.get("encoding", "utf-8") or "utf-8"

    try:
        return raw.decode(encoding)
    except (UnicodeDecodeError, LookupError):
        return raw.decode("utf-8", errors="replace")


# ─── PDF 解析（增强）────────────────────────────────────────

def _parse_pdf(file_path: str) -> str:
    """
    解析 PDF 文件，优先提取文字；无文字时提示 OCR

    表格处理策略（基于 PDF 结构分析）:
    1. PyPDF2 提取文字 + 位置信息
    2. 检测近似等宽对齐的文本块（表格假设）
    3. 合并为 Markdown 表格格式输出
    """
    try:
        from PyPDF2 import PdfReader
    except ImportError:
        raise ImportError(
            "PDF 解析需要 PyPDF2 库: pip install PyPDF2"
        )

    reader = PdfReader(file_path)
    if reader.is_encrypted:
        try:
            reader.decrypt("")
        except Exception:
            pass  # 尝试不解密读取

    pages_result = []
    total_pages = len(reader.pages)

    for page_num, page in enumerate(reader.pages):
        text = page.extract_text()

        if not text or not text.strip():
            # 扫描 PDF：无可提取文字
            pages_result.append(
                f"[第 {page_num + 1} 页 — 此页无可提取文字（扫描件或图片页）]"
            )
            continue

        # 尝试表格检测
        tables = _extract_tables_from_pdf_text(text)

        if tables:
            # 有表格：每个表格单独输出
            for table_md in tables:
                pages_result.append(table_md)
            # 保留表格间的非表格文字
            non_table = _remove_tables_from_text(text)
            if non_table.strip():
                pages_result.append(non_table.strip())
        else:
            pages_result.append(text.strip())

    combined = "\n\n".join(pages_result)

    if not combined.strip():
        raise ValueError(
            f"PDF 文件无可提取文字内容（扫描件）: {file_path}\n"
            "建议：使用支持 OCR 的工具（如 pdf2image + pytesseract）先将扫描件转为文字"
        )

    # 在开头标注来源
    source_note = f"[来源: {os.path.basename(file_path)}, 共 {total_pages} 页]\n"
    return source_note + combined


def _extract_tables_from_pdf_text(text: str) -> list[str]:
    """
    从 PDF 提取的文字块中检测表格结构
    策略：等宽空格对齐行 → 检测分隔符 → 转为 Markdown 表格

    返回: Markdown 格式的表格列表
    """
    lines = text.split("\n")
    table_candidates = []
    i = 0

    while i < len(lines):
        line = lines[i]
        # 检测是否为表格行：含 2 个以上 | 分隔符 或 3+ 个空格对齐的列
        pipe_count = line.count("|")
        spaces = line.split(" ")

        is_table_row = (
            pipe_count >= 3  # Markdown 风格表格
            or (len(spaces) >= 4 and _looks_like_table_row(line))  # 空格等宽表格
        )

        if is_table_row:
            # 收集连续表格行
            table_lines = [line]
            j = i + 1
            while j < len(lines):
                next_line = lines[j]
                # 表格行通常有相似的分隔符密度
                if next_line.strip() and (
                    next_line.count("|") >= 2
                    or _looks_like_table_row(next_line)
                    or next_line.strip() == ""
                ):
                    if next_line.strip() and next_line.strip() != "":
                        table_lines.append(next_line)
                    j += 1
                else:
                    break

            # 尝试转换
            md_table = _convert_to_markdown_table(table_lines)
            if md_table:
                table_candidates.append(md_table)

            i = j
        else:
            i += 1

    return table_candidates


def _looks_like_table_row(line: str) -> bool:
    """判断一行是否为表格数据行（基于空格对齐分析）"""
    stripped = line.strip()
    if not stripped or stripped.startswith("#") or "│" in line:
        return False

    # 分析空格分布：表格列通常有规律的空格分隔
    parts = stripped.split()
    if len(parts) < 2:
        return False

    # 计算相邻词之间的空格数方差
    char_pos = []
    pos = 0
    for part in parts:
        p = stripped.find(part, pos)
        char_pos.append(p)
        pos = p + len(part)

    if len(char_pos) < 2:
        return False

    gaps = [char_pos[k] - char_pos[k - 1] for k in range(1, len(char_pos))]
    if not gaps:
        return False

    avg_gap = sum(gaps) / len(gaps)
    variance = sum((g - avg_gap) ** 2 for g in gaps) / len(gaps)

    # 表格列间距通常相对均匀（方差较小）
    return variance < avg_gap * 2 and avg_gap >= 2


def _convert_to_markdown_table(table_lines: list[str]) -> str | None:
    """
    将原始文本行转换为 Markdown 表格
    策略：
    1. 识别表头（第一行）
    2. 识别分隔行（---, ===, |---|）
    3. 解析内容行
    """
    if not table_lines:
        return None

    # 过滤空行
    valid_lines = [l for l in table_lines if l.strip()]
    if len(valid_lines) < 2:
        return None

    md_lines = []

    for line in valid_lines:
        line = line.strip()

        # Markdown 格式表格行
        if line.startswith("|"):
            cells = [c.strip() for c in line.strip("|").split("|")]
            # 过滤分隔行
            if all(re.match(r"^[-:\s]+$", c) for c in cells):
                md_lines.append(line)
            else:
                md_lines.append(line)
        else:
            # 非 Markdown 格式：转为 |
            cells = re.split(r"\s{2,}", line)
            if len(cells) >= 2:
                md_line = "| " + " | ".join(c.strip() for c in cells if c.strip()) + " |"
                md_lines.append(md_line)

    if not md_lines:
        return None

    # 确保有分隔行
    if not any("---" in l or "--" in l for l in md_lines):
        # 找到第一个可能是表头的行，在其后插入分隔行
        for idx, l in enumerate(md_lines):
            cells = [c.strip() for c in l.strip("|").split("|") if c.strip()]
            if len(cells) >= 2:
                sep = "| " + " | ".join(["---"] * len(cells)) + " |"
                md_lines.insert(idx + 1, sep)
                break

    result = "\n".join(md_lines)
    # 简单验证：至少有表头和分隔行
    if md_lines.count("|") >= 4:
        return result

    return None


def _remove_tables_from_text(text: str) -> str:
    """从文本中移除已提取的表格部分，保留普通文字"""
    lines = text.split("\n")
    cleaned = []
    skip_next = 0

    for line in lines:
        if skip_next > 0:
            skip_next -= 1
            continue

        # 检测是否为表格分隔行（全是 - : | 空格）
        if re.match(r"^\|?[\s\-:|#]+\|?[\s\-:|#]*$", line.strip()):
            continue

        pipe_count = line.count("|")
        if pipe_count >= 4:  # 可能是表格行
            # 检查是否在连续表格块中
            cleaned.append("")  # 留空，后续可以清理
        else:
            if line.strip():
                cleaned.append(line)

    return "\n".join(cleaned)


# ─── DOCX 解析（增强）───────────────────────────────────────

def _parse_docx(file_path: str) -> str:
    """
    解析 DOCX 文件，提取段落、标题、表格

    表格处理：
    - 合并单元格 → 保留原始文本
    - 嵌套表格 → 扁平化处理
    - 输出为 Markdown 表格格式
    """
    try:
        from docx import Document
        from docx.enum.text import WD_PARAGRAPH_TYPE
    except ImportError:
        raise ImportError(
            "DOCX 解析需要 python-docx: pip install python-docx"
        )

    doc = Document(file_path)
    blocks = []

    # 遍历所有元素（段落 + 表格交替）
    for element in doc.element.body:
        tag = element.tag.split("}")[-1]

        if tag == "p":
            # 段落
            para = element
            # 获取段落文本
            text = "".join(node.text or "" for node in para.iter())
            if text.strip():
                # 检测是否为标题样式（简略：开头有 # 或特殊格式）
                if text.startswith("#") or re.match(r"^(第[一二三四五六七八九十]+条|第\d+章)", text):
                    blocks.append(text)
                else:
                    blocks.append(text)

        elif tag == "tbl":
            # 表格元素 → 转为 Markdown
            tbl_md = _docx_table_to_markdown(doc, element)
            if tbl_md:
                blocks.append(tbl_md)

    # 合并段落文字
    result_parts = []
    for block in blocks:
        if block.strip():
            result_parts.append(block.strip())

    # 单独提取段落中的表格
    for table in doc.tables:
        tbl_md = _docx_table_to_markdown_from_obj(table)
        if tbl_md:
            result_parts.append(tbl_md)

    combined = "\n\n".join(result_parts)
    if not combined.strip():
        raise ValueError(f"DOCX 文件无文本内容: {file_path}")

    return f"[来源: {os.path.basename(file_path)}]\n" + combined


def _docx_table_to_markdown(doc, table_element) -> str | None:
    """将 DOCX table 元素转为 Markdown"""
    try:
        from docx.table import Table
        from docx.text.paragraph import Paragraph

        tbl = Table(table_element, doc)
        rows = list(tbl.rows)
        if len(rows) < 2:
            return None

        lines = []
        for row in rows:
            cells = [cell.text.strip() for cell in row.cells]
            lines.append("| " + " | ".join(cells) + " |")

        # 插入分隔行
        if len(lines) >= 1:
            col_count = len(lines[0].split("|")[:-1])
            sep = "| " + " | ".join(["---"] * (col_count - 1)) + " |"
            lines.insert(1, sep)

        return "\n".join(lines) if lines else None
    except Exception:
        return None


def _docx_table_to_markdown_from_obj(table) -> str | None:
    """将 python-docx Table 对象转为 Markdown"""
    rows = list(table.rows)
    if len(rows) < 2:
        return None

    lines = []
    for row in rows:
        cells = []
        for cell in row.cells:
            # 处理合并单元格
            text = cell.text.strip().replace("\n", " ")
            cells.append(text)
        lines.append("| " + " | ".join(cells) + " |")

    if lines:
        col_count = len(lines[0].split("|")[:-1])
        sep = "| " + " | ".join(["---"] * (col_count - 1)) + " |"
        lines.insert(1, sep)
        return "\n".join(lines)
    return None


# ─── XLSX 解析 ─────────────────────────────────────────────

def _parse_xlsx(file_path: str) -> str:
    """
    解析 XLSX 文件，每个 Sheet 输出一节
    表格输出为 Markdown 格式
    """
    try:
        import openpyxl
    except ImportError:
        raise ImportError(
            "XLSX 解析需要 openpyxl: pip install openpyxl"
        )

    wb = openpyxl.load_workbook(file_path, data_only=True)
    sections = []
    filename = os.path.basename(file_path)

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows_data = []

        for row in ws.iter_rows(values_only=True):
            # 过滤全空行
            if any(cell is not None for cell in row):
                rows_data.append(row)

        if not rows_data:
            continue

        # 转换每个 Sheet 为 Markdown 表格
        lines = []
        for ri, row in enumerate(rows_data):
            cells = [str(cell) if cell is not None else "" for cell in row]
            lines.append("| " + " | ".join(cells) + " |")

            # 表头后加分隔行
            if ri == 0:
                col_count = len(cells)
                lines.append("| " + " | ".join(["---"] * col_count) + " |")

        section = f"## Sheet: {sheet_name}\n\n" + "\n".join(lines)
        sections.append(section)

    if not sections:
        raise ValueError(f"XLSX 文件为空: {file_path}")

    return f"[来源: {filename}]\n\n" + "\n\n".join(sections)


# ─── PPTX 解析 ─────────────────────────────────────────────

def _parse_pptx(file_path: str) -> str:
    """
    解析 PPTX 文件，每张幻灯片输出一段
    提取标题 + 正文文本
    """
    try:
        import pptx
        from pptx.util import Pt
    except ImportError:
        raise ImportError(
            "PPTX 解析需要 python-pptx: pip install python-pptx"
        )

    prs = pptx.Presentation(file_path)
    slides = []
    filename = os.path.basename(file_path)

    for slide_num, slide in enumerate(prs.slides, 1):
        texts = []

        # 提取所有形状中的文字
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    texts.append(text)

        # 提取表格
        for shape in slide.shapes:
            if shape.has_table:
                tbl = shape.table
                rows = [[cell.text.strip() for cell in row.cells] for row in tbl.rows]
                for row in rows:
                    texts.append("| " + " | ".join(row) + " |")
                # 表头分隔行
                col_count = len(rows[0]) if rows else 0
                if col_count > 0:
                    texts.append("| " + " | ".join(["---"] * col_count) + " |")

        if texts:
            slide_text = "\n".join(texts)
            slides.append(f"### 幻灯片 {slide_num}\n\n{slide_text}")

    if not slides:
        raise ValueError(f"PPTX 文件无文本内容: {file_path}")

    return f"[来源: {filename}]\n\n" + "\n\n---\n\n".join(slides)


# ─── 导入到知识库 ───────────────────────────────────────────

def import_file_to_kb(
    file_path: str,
    kb_name: str,
    tags: list = None,
    memory_manager=None,
) -> dict:
    """导入文档到知识库（支持所有已解析格式）"""
    if not os.path.exists(file_path):
        return {"error": f"文件不存在: {file_path}"}

    try:
        text = parse_file(file_path)
    except Exception as e:
        return {"error": f"解析失败: {e}"}

    if not text.strip():
        return {"error": "文件内容为空"}

    if memory_manager is None:
        from memory_manager import MemoryManager
        memory_manager = MemoryManager()

    filename = os.path.basename(file_path)
    all_tags = list(tags or [])
    all_tags.append(f"file:{filename}")

    result = memory_manager.add_to_kb(
        kb_name=kb_name,
        text=text,
        source_file=filename,
        tags=all_tags,
    )

    return {
        "status": "imported",
        "file": filename,
        "format": os.path.splitext(file_path)[1].lower(),
        "text_length": len(text),
        "result": result,
    }


def import_directory_to_kb(
    dir_path: str,
    kb_name: str,
    tags: list = None,
    memory_manager=None,
) -> dict:
    """
    批量导入目录下的文档到知识库

    支持格式：.txt, .md, .pdf, .docx, .xlsx, .pptx
    """
    if not os.path.isdir(dir_path):
        return {"error": f"目录不存在: {dir_path}"}

    supported_ext = {
        ".txt", ".md", ".markdown",
        ".pdf",
        ".docx", ".doc",
        ".xlsx",
        ".pptx",
    }

    files = []
    for root, _, filenames in os.walk(dir_path):
        for fname in filenames:
            ext = os.path.splitext(fname)[1].lower()
            if ext in supported_ext:
                files.append(os.path.join(root, fname))

    if not files:
        return {
            "error": (
                f"目录中没有支持的文件（支持: "
                + ", ".join(sorted(supported_ext)) + "）"
            )
        }

    if memory_manager is None:
        from memory_manager import MemoryManager
        memory_manager = MemoryManager()

    results = []
    success, failed = 0, 0

    for fpath in files:
        r = import_file_to_kb(fpath, kb_name, tags, memory_manager)
        results.append(r)
        if r.get("status") == "imported":
            success += 1
        else:
            failed += 1

    return {
        "status": "batch_imported",
        "total_files": len(files),
        "success": success,
        "failed": failed,
        "details": results,
    }


# ─── CLI 入口 ───────────────────────────────────────────────
if __name__ == "__main__":
    import sys

    if len(sys.argv) < 2:
        print("用法: python doc_parser.py <file_path>")
        sys.exit(1)

    filepath = sys.argv[1]
    try:
        text = parse_file(filepath)
        print(text[:2000])  # 最多打印前 2000 字
        if len(text) > 2000:
            print(f"\n... [省略 {len(text) - 2000} 字] ...")
    except Exception as e:
        print(f"错误: {e}", file=sys.stderr)
        sys.exit(1)
